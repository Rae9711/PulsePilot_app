from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

NAVY = RGBColor(21, 43, 77)
TEAL = RGBColor(28, 153, 135)
LIGHT = RGBColor(245, 249, 252)
DARK = RGBColor(30, 30, 30)
WHITE = RGBColor(255, 255, 255)


def style_title(title_shape, subtitle_shape=None):
    tf = title_shape.text_frame
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if subtitle_shape is not None:
        stf = subtitle_shape.text_frame
        for para in stf.paragraphs:
            para.font.size = Pt(20)
            para.font.color.rgb = DARK


def add_header_bar(slide, text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tx = bar.text_frame
    tx.clear()
    p = tx.paragraphs[0]
    p.text = text
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT


def add_bullets(title, bullets):
    s = prs.slides.add_slide(prs.slide_layouts[1])
    add_header_bar(s, "FitForecast - Final Presentation")
    s.shapes.title.text = title
    s.shapes.title.left = Inches(0.6)
    s.shapes.title.top = Inches(0.9)
    s.shapes.title.width = Inches(12.0)
    s.shapes.title.height = Inches(0.8)
    s.shapes.title.text_frame.paragraphs[0].font.size = Pt(34)
    s.shapes.title.text_frame.paragraphs[0].font.bold = True
    s.shapes.title.text_frame.paragraphs[0].font.color.rgb = NAVY

    body = s.shapes.placeholders[1]
    body.left = Inches(0.9)
    body.top = Inches(1.9)
    body.width = Inches(11.8)
    body.height = Inches(5.0)

    tf = body.text_frame
    tf.clear()
    first = True
    for item in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p.text = text
        p.level = level
        p.font.size = Pt(24 if level == 0 else 20)
        p.font.color.rgb = DARK
    return s


# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[0])
add_header_bar(slide, "FitForecast - Final Presentation")
slide.shapes.title.text = "FitForecast"
slide.placeholders[1].text = "Discovering personal fitness patterns through explainable insights"
style_title(slide.shapes.title, slide.placeholders[1])

# Slide 2
add_bullets("The Problem (Relatable)", [
    "Most fitness apps give generic advice that ignores individual context.",
    "People ask: Why do I feel amazing after some sessions and drained after others?",
    "Without personal pattern tracking, users cannot tell what actually works for them.",
    "Result: inconsistent routines, frustration, and stalled progress.",
])

# Slide 3
add_bullets("Our Solution", [
    "FitForecast is a personalization-first fitness pattern analyzer.",
    "Users log workouts/meals in natural language plus pre/post feelings (mood, energy, stress).",
    "System learns each user's baseline and surfaces explainable insights.",
    "Predictions compare scenarios so users can test timing, fueling, and routine changes.",
])

# Slide 4
add_bullets("Why It Matters", [
    "Moves users from guesswork to data-informed self-discovery.",
    "Encourages sustainable behavior change with personalized feedback.",
    "Works for different user types: improver, inconsistent/high-stress, peak performer.",
    "Insights remain useful even when optional LLM wording is offline.",
])

# Slide 5
add_bullets("What We Built", [
    "Full-stack app: React + TypeScript frontend, Node/Express + Prisma backend, PostgreSQL.",
    "JWT authentication, protected routes, and per-user data isolation.",
    "Dashboard, logging flow, history, trends, and predictions.",
    "Seeded demo personas with realistic six-month histories (95+ entries each).",
    "Validation stream with analytics notebooks and rule testing.",
])

# Slide 6
s = add_bullets("Landing Page Walkthrough", [
    "Show login/landing experience and clear value proposition.",
    "Point out demo persona quick-login options (Athena, Boris, Cora).",
    "Highlight navigation: Dashboard, Log, History, Trends.",
    "Explain how users move from entry logging to insights and predictions.",
])
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(2.1), Inches(5.7), Inches(4.2))
box.fill.solid()
box.fill.fore_color.rgb = LIGHT
box.line.color.rgb = TEAL
bx = box.text_frame
bx.clear()
p = bx.paragraphs[0]
p.text = "Insert landing page screenshot here"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY
p.alignment = PP_ALIGN.CENTER

# Slide 7
add_bullets("Live Demo Flow", [
    "1) Login as Athena (consistent improver).",
    "2) Dashboard: show insights and prediction panel.",
    "3) History: verify learned patterns from real entries.",
    "4) Trends: show baselines and changes over time.",
    "5) Log a new workout plus feelings and show refresh impact.",
    "6) Optional: switch to Boris to show contrast in recommendations.",
])

# Slide 8
add_bullets("Architecture Snapshot", [
    "Stream 1: Backend API plus business logic plus database.",
    "Stream 2: Frontend experience and user workflows.",
    "Stream 3: Analytics notebooks for rule validation.",
    "Stream 4: Integration scenarios plus optional LLM narrative rewriting.",
    "Design principle: deterministic rules for trust, optional AI for wording polish.",
])

# Slide 9
add_bullets("Reflection: Goal vs Outcome", [
    "Beginning-of-semester goal:",
    ("Build an app that helps people understand how habits affect energy and mood.", 1),
    "What we actually built:",
    ("A complete MVP with auth, personalized baselines, insights, trends, and forecasts.", 1),
    ("A robust demo environment with multiple personas and realistic longitudinal data.", 1),
    "Key learning:",
    ("Personalization plus explainability creates stronger user trust than one-size-fits-all advice.", 1),
])

# Slide 10
add_bullets("Challenges and Wins", [
    "Challenge: making insights both personal and explainable.",
    "Win: deterministic signal pipeline with user-specific baselines.",
    "Challenge: proving value in a short semester timeline.",
    "Win: end-to-end product flow and demo-ready personas with rich data.",
    "Challenge: balancing AI with reliability.",
    "Win: app remains fully functional without LLM dependency.",
])

# Slide 11
add_bullets("Submission Requirement: Video Checklist", [
    "Required clip 1: walkthrough of your landing page.",
    "Required clip 2: demo of your application end-to-end.",
    "Recommended length: 3-6 minutes.",
    "Suggested structure: 45s problem + 60s solution + 2-3 min demo + 45s reflection.",
    "Export and submit alongside this slide deck.",
])

# Slide 12
add_bullets("Thank You", [
    "FitForecast helps users discover what works for their body and routine.",
    "Questions?",
])

OUTPUT = "FitForecast_Final_Presentation.pptx"
prs.save(OUTPUT)
print(OUTPUT)
